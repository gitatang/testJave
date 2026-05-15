# -*- coding: utf-8 -*-
"""
2025年年终总结PPT生成器
北京艾棋_M9_技术部_汤松杰
基于Q1-Q4季度考核真实内容生成
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("正在安装 python-pptx...")
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "-q"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle=""):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True

        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(22)
        subtitle_para.font.color.rgb = RGBColor(68, 68, 68)


def add_content_slide(prs, title, content_items, show_number=True, font_size=20):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        if show_number:
            p.text = f"{i+1}. {item}"
        else:
            p.text = f"• {item}"

        p.font.size = Pt(font_size)
        p.space_after = Pt(10)
        p.level = 0


def create_ppt():
    """创建年终总结PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== 第1页：封面 ==========
    add_title_slide(prs, "2025年度年终总结",
                   "北京艾棋_M9_技术部_汤松杰\n2026年2月")

    # ========== 第2页：目录 ==========
    add_content_slide(prs, "目录", [
        "年度工作概述",
        "Q1（1-3月）基础服务建设",
        "Q2（4-6月）核心模块开发",
        "Q3（7-9月）服务优化与集成",
        "Q4（10-12月）监控体系建设",
        "年度成长与进步",
        "待提高项",
        "2026年度规划"
    ], show_number=False, font_size=24)

    # ========== 第3页：年度工作概述 ==========
    add_content_slide(prs, "年度工作概述", [
        "【核心服务】完成M9项目多个核心服务的开发与维护",
        "  - Push服务、Rank服、广播服、背包服、任务服、工具服、红点服",
        "",
        "【系统集成】完成A8平台短信、邮件发送能力集成",
        "",
        "【新模块开发】M9俱乐部服务、VIP任务模块从0到1开发完成",
        "",
        "【监控体系】构建M9项目监控服务体系，完成上报接口与多渠道告警",
        "",
        "【技术成长】形成\"分析-设计-开发\"标准化工作流程"
    ], font_size=19)

    # ========== 第4页：Q1工作总览 ==========
    add_content_slide(prs, "Q1（1-3月）工作总览", [
        "【Push服务】邮箱验证码发送功能",
        "  - 基于RocketMQ监听消息发送邮件",
        "  - 支持多语言自定义邮箱发送",
        "",
        "【Rank服】排行榜服务开发",
        "  - 多参数综合排序（段位、星级、积分）",
        "  - 支持排行类型扩展和分区",
        "  - Redis+MySQL双重存储，支持热更",
        "",
        "【广播服】公告模块开发",
        "  - 后台配置、定时刷新、多语言支持",
        "  - 红点提示机制"
    ], font_size=18)

    # ========== 第5页：Q1详细工作 - Rank服 ==========
    add_content_slide(prs, "Q1详细：Rank服排行榜服务", [
        "✅ 多参数排序：支持按段位、星级、积分综合排序",
        "✅ 类型扩展：可通过添加配置和表的方式扩展多种排行榜",
        "✅ 分区规则：编写快速分区规则，将玩家快速完成分区",
        "✅ 数据落地：排行榜数据落地到MySQL",
        "✅ 容灾机制：Redis崩溃或数据丢失时，可通过热更方式重新加载数据，无需重启"
    ], font_size=19)

    # ========== 第6页：Q1详细工作 - 广播服公告模块 ==========
    add_content_slide(prs, "Q1详细：广播服公告模块", [
        "✅ 开发公告模块：设计公告配置表，支持后台添加后定时刷新到服务器",
        "✅ 多语言支持：公告支持多语言选择",
        "✅ 自定义排序：支持公告自定义排序规则",
        "✅ 读取记录：实现玩家读取公告后保存记录",
        "✅ 红点提示：支持前端显示未读取的红点提示"
    ], font_size=19)

    # ========== 第7页：Q2工作总览 ==========
    add_content_slide(prs, "Q2（4-6月）工作总览", [
        "【背包服】商城与背包系统",
        "  - 商城物品配置、购买兑换、背包查询",
        "  - 真实充值物品更新、订单退款处理",
        "",
        "【任务服】任务、成就、通行证系统",
        "  - 任务列表、成就系统、通行证完整开发",
        "  - MQ事件驱动、多语言支持、充值激活",
        "",
        "【工具服】实用组件开发",
        "  - 胜率计算器、记分牌组件",
        "",
        "【红点服】红点服务接入"
    ], font_size=18)

    # ========== 第8页：Q2详细 - 背包服 ==========
    add_content_slide(prs, "Q2详细：背包服商城系统", [
        "✅ 商城配置：支持添加删除商城物品，支持多语言",
        "✅ 购买限制：商品支持限制购买数量，支持多种刷新方式（不刷新/每日/每周/每月）",
        "✅ 上架下架：支持上架下架删除和开放时间限制",
        "✅ 物品有效期：可以配置物品有效期",
        "✅ 金块商店：支持金块（一级货币）购买",
        "✅ 充值接口：提供真实充值后背包物品更新接口",
        "✅ 订单处理：判断同一订单是否多次发放，支持交易退款返还"
    ], font_size=17)

    # ========== 第9页：Q2详细 - 任务服 ==========
    add_content_slide(prs, "Q2详细：任务服任务与通行证", [
        "✅ 任务系统：支持多语言、前置任务、刷新类型、开始结束时间",
        "✅ 成就系统：配置化管理，支持上线下架、时间控制",
        "✅ MQ事件：完成事件MQ，分析消息添加任务/成就进度，更新记录最大值",
        "✅ 通行证系统：完整通行证功能，支持充值激活、购买等级",
        "✅ 奖励领取：完成任务/成就/通行证可领取对应奖励",
        "✅ 周期重启：通行证结束后可重新开始，重新领奖"
    ], font_size=18)

    # ========== 第10页：Q2详细 - 工具服 ==========
    add_content_slide(prs, "Q2详细：工具服实用组件", [
        "✅ 胜率计算器：单机版胜率计算组件",
        "  - 玩家可自主添加底牌和手牌计算胜率",
        "  - 支持选择器计算范围胜率和名次",
        "",
        "✅ 记分牌组件：对局筹码变化与数据统计",
        "  - 添加玩家、创建对局、模拟初始筹码",
        "  - 统计进出筹码，计算总计筹码",
        "  - 统计局数、查看详情、时间查询、操作记录",
        "",
        "✅ 红点服务：公告、任务、成就、通行证模块红点推送"
    ], font_size=18)

    # ========== 第11页：Q3工作成果 ==========
    add_content_slide(prs, "Q3（7-9月）工作成果", [
        "✅ M9俱乐部服务：完成开发、测试及生产环境部署，服务运行平稳",
        "",
        "✅ M9VIP任务模块：顺利完成开发、测试及部署全流程，正式投入运营",
        "",
        "✅ Bug修复：测试阶段反馈的Bug全部修复并验证，代码已更新",
        "",
        "✅ Push服务集成：成功集成A8的短信与邮件发送能力，功能开发联调完毕",
        "",
        "📈 工作方法转变：形成\"分析-设计-开发\"标准化流程",
        "  - 首先深入分析需求本质并评估扩展性",
        "  - 然后进行服务抽象与接口设计",
        "  - 最后进入开发阶段"
    ], font_size=18)

    # ========== 第12页：Q4工作成果 ==========
    add_content_slide(prs, "Q4（10-12月）工作成果", [
        "【监控服务】从0到1构建监控服务系统",
        "  ✅ 完成上报接口开发，其他服务可上报错误信息",
        "  ✅ 支持单个或多个接收人配置",
        "  ✅ 集成6种通知渠道：短信、邮箱、钉钉、飞书、Lark、Telegram",
        "  ✅ 根据配置将告警消息发送给指定负责人",
        "",
        "【俱乐部服务优化】7项功能新增",
        "  ✅ 模糊查询俱乐部及成员",
        "  ✅ 自定义俱乐部排序规则",
        "  ✅ 配置用户创建俱乐部数量上限",
        "  ✅ 俱乐部成员自定义排序",
        "  ✅ 自动审批开关",
        "  ✅ 红点提示及刷新机制优化",
        "",
        "【质量保障】修复线上已知Bug，跟进处理Bug清单"
    ], font_size=17)

    # ========== 第13页：年度成长与进步 ==========
    add_content_slide(prs, "年度成长与进步", [
        "【思维模式转变】",
        "  - 从\"接到需求即编码\"转变为\"分析-设计-开发\"标准化流程",
        "  - 解构需求本质，评估长期扩展性",
        "  - 进行服务抽象与接口契约设计，明确系统边界",
        "",
        "【沟通协作提升】",
        "  - 养成对模糊需求主动澄清的习惯",
        "  - 优先向相关同事或产品经理请教，确保理解一致",
        "  - 有效避免因信息偏差导致的返工",
        "",
        "【工具应用探索】",
        "  - 开始探索AI工具在开发工作流中的应用",
        "  - 技术方案咨询、单元测试/样板代码生成、日志分析等场景"
    ], font_size=18)

    # ========== 第14页：待提高项 ==========
    add_content_slide(prs, "待提高项", [
        "【技术深度】",
        "  - 边界情况和极限情况考虑不够清晰",
        "  - 需要更深入思考技术方案的简洁性和有效性",
        "",
        "【工具使用平衡】",
        "  - 曾经过度依赖AI生成代码，未充分批判性评估",
        "  - 导致采用过度设计、不符合\"简单有效\"原则的方案",
        "  - 认识到工具价值在于拓展思路，而非替代思考",
        "",
        "【创新精神】",
        "  - 面对常规任务时，主动挑战的锐气和探索精神不足",
        "  - 需要从\"完成开发\"转向\"创造价值\"",
        "",
        "【需求理解】",
        "  - 项目需求需要全面了解和深挖",
        "  - 多数精力聚焦于\"交付\"，未深入思考迭代和优化"
    ], font_size=18)

    # ========== 第15页：2026年度规划 ==========
    add_content_slide(prs, "2026年度规划", [
        "【AI应用深化】",
        "  - 将AI应用场景化、习惯化",
        "  - 重点场景：技术方案咨询、单元测试生成、日志分析",
        "  - 季度末复盘提效效果",
        "",
        "【思维转变】",
        "  - 从\"完成开发\"转向\"创造价值\"",
        "  - 通过深度参与提升工作投入感和成就感",
        "",
        "【技术方案优化】",
        "  - 在技术方案决策与工具辅助之间找到平衡",
        "  - 对AI输出进行批判性评估与简化重构",
        "  - 坚持\"简单有效\"原则，避免过度设计",
        "",
        "【持续改进】",
        "  - 深入思考怎么更便于迭代和优化",
        "  - 不止于交付，更要关注长期价值"
    ], font_size=18)

    # ========== 第16页：结束页 ==========
    add_title_slide(prs, "感谢聆听", "期待2026年再创佳绩！\n从\"完成\"到\"创造\"，持续成长")

    # 保存文件
    output_path = r"D:\tangsongjie\work\github\TestJave\北京艾棋_M9_技术部_汤松杰_2025年终总结_完整版.pptx"
    prs.save(output_path)
    print("=" * 70)
    print("PPT generation completed!")
    print(f"Output: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print("Content: Based on Q1-Q4 quarterly performance reviews")
    print("=" * 70)
    return output_path


if __name__ == "__main__":
    create_ppt()